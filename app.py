import os
import google.generativeai as genai
import streamlit as st
from pptx import Presentation
from pptx.util import Inches  # Fixed: Missing Import
import json
import re

# --- CONFIGURATION ---
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
genai.configure(api_key=API_KEY)
# Updated to Gemini 3 for 2026 stability
model = genai.GenerativeModel('gemini-3.1-flash-lite-preview')


API_KEY = "AIzaSyBQvch_1pEp0FHZ0t0VP32sJbwa-4Wlf8w"


# --- SIDEBAR (Student Info) ---
with st.sidebar:
    st.title("🎓 Project Credits")
    student_name = st.text_input("Name", placeholder="Your Name")
    student_usn = st.text_input("USN", placeholder="Your USN")
    st.info("VTU Phase 1: AI PPT Generator")

# --- MAIN INTERFACE ---
st.title("🎨 AI Automatic PPT Creator")

user_input = st.text_area("Paste Content:", height=200)
num_slides = st.slider("Slides", 3, 10, 5)

if st.button("✨ Generate Presentation"):
    if not user_input or not student_name:
        st.error("Please fill in Name and Content!")
    else:
        with st.spinner("AI is working..."):
            try:
                # 1. GENERATE DATA
                prompt = f"Summarize into {num_slides} slides in JSON: [{{'title':'T', 'content':['P1']}}]. Text: {user_input}"
                response = model.generate_content(prompt)
                
                # 2. EXTRACT JSON (This defines slides_data)
                raw_json = re.search(r'\[.*\]', response.text, re.DOTALL).group()
                slides_data = json.loads(raw_json)
                
                # 3. CREATE PPT (This defines prs)
                prs = Presentation()
                
                # Title Slide
                title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                # This takes the first 50 characters of your input as the title
                title_slide.shapes.title.text = user_input[:50].title()
                title_slide.placeholders[1].text = f"By: {student_name}\nUSN: {student_usn}"
                
                # 1. Set the Main Topic Title (Dynamic)
                title_slide.shapes.title.text = user_input.split('\n')[0][:60].title()
                
                # 2. Set the Subtitle Text (Name and USN)
                subtitle_placeholder = title_slide.placeholders[1]
                
                subtitle_placeholder.text = f"Presented by: {student_name}\nUSN: {student_usn}"
                
                # 3. ADD THE FORMATTING HERE
                from pptx.util import Pt # Make sure this is at the top of your file or her
                subtitle_frame = subtitle_placeholder.text_frame
                p = subtitle_frame.paragraphs[0]
                p.font.size = Pt(24)
                p.font.bold = True
                # Content Slides (Using slides_data)
                
                for item in slides_data:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = item['title']
                    slide.placeholders[1].text = "\n".join(item['content'])
                
                # 4. SAVE AND DOWNLOAD
                file_path = "output.pptx"
                prs.save(file_path)
                
                st.success("Done!")
                with open(file_path, "rb") as f:
                    st.download_button("📥 Download PPT", f, file_name="Project.pptx")
                    
            except Exception as e:
                st.error(f"Logic Error: {e}")